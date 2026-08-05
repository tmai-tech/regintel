#!/usr/bin/env python3
"""
Ministry Document Crawler
=========================
Crawls a ministry (government) website, discovers downloadable documents
(PDF / DOC / DOCX / XLS / XLSX / PPT / PPTX / CSV / RTF / ODT / ODS / TXT),
downloads them, extracts text, detects the language, and emits a machine-
readable index that an agent can use to translate + summarize each document.

Built to survive real-world government-site obstacles:
  - Bot protection / WAF (Cloudflare, Akamai): realistic browser headers,
    session cookies, exponential-backoff retries, block-page detection.
  - Rate limiting / polite crawling: configurable delay, page caps.
  - JS-rendered SPAs: scans raw HTML *and* inline JSON/script blobs for
    document URLs; if under-discovery is suspected the caller can re-run
    against the Firecrawl integration (JS rendering) as a fallback.
  - Duplicate links: URL normalization + visited set.
  - Huge files: size cap with skip-and-note.
  - Scanned PDFs: detected (no extractable text) and flagged for OCR.
  - Mixed encodings: UTF-8 -> cp1252 -> latin-1 fallback chain.
  - Geo/VPN blocks: unreachable hosts are reported clearly, not fabricated.

Usage:
  python3 ministry_crawler.py --url https://www.example.gov \
      --output-dir ./crawl_out [options]

Options:
  --url URL                 Ministry site root (required)
  --output-dir DIR          Where docs/, text/, index.json are written (default: ./crawl_out)
  --max-pages N             Max pages to crawl (default: 300)
  --delay SECONDS           Delay between HTTP requests (default: 0.8)
  --max-file-mb N           Skip downloads larger than N MB (default: 150)
  --max-text-chars N        Cap extracted text per doc (default: 20000)
  --include-other-domains   Also follow links to non-ministry domains (CDNs etc.)
  --enforce-robots          Respect robots.txt (default: off - gov sites often block bots)
  --translate               Best-effort machine translation of non-English text
                            (deep-translator -> Google; may be rate-limited). The
                            agent's LLM translation is preferred; this is a fallback.
  --enable-ocr              OCR scanned/image-based PDFs with tesserocr (requires pdf2image + tesserocr)
  --insecure                Skip TLS verification (only for broken-proxy sandboxes)
  --seed-pages URL [URL...] Extra start pages (e.g. /publications, /documents)
"""

import argparse, csv, json, os, re, sys, time, zipfile
from urllib.parse import urljoin, urlparse, urlunparse, unquote
import requests
from bs4 import BeautifulSoup

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except Exception:
    HAS_FITZ = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except Exception:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook
    HAS_XLSX = True
except Exception:
    HAS_XLSX = False

try:
    import xlrd
    HAS_XLS = True
except Exception:
    HAS_XLS = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

try:
    from py3langid.langid import LanguageIdentifier, MODEL_FILE
    _lid = LanguageIdentifier.from_pickled_model(MODEL_FILE, norm_probs=True)
    HAS_LANGID = True
except Exception:
    HAS_LANGID = False

try:
    from deep_translator import GoogleTranslator
    HAS_TRANSLATOR = True
except Exception:
    HAS_TRANSLATOR = False

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except Exception:
    HAS_PDF2IMAGE = False

try:
    from tesserocr import PyTessBaseAPI, PSM
    HAS_TESSEROCR = True
    import os as _os
    _TESSDATA = _os.environ.get('TESSDATA_PREFIX', '/tmp/tessdata')
except Exception:
    HAS_TESSEROCR = False

DOC_EXTENSIONS = {
    '.pdf': 'PDF', '.doc': 'DOC', '.docx': 'DOCX', '.xls': 'XLS',
    '.xlsx': 'XLSX', '.ppt': 'PPT', '.pptx': 'PPTX', '.csv': 'CSV',
    '.rtf': 'RTF', '.odt': 'ODT', '.ods': 'ODS', '.txt': 'TXT',
}
CONTENT_TYPE_MAP = {
    'application/pdf': 'PDF',
    'application/msword': 'DOC',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'DOCX',
    'application/vnd.ms-excel': 'XLS',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'XLSX',
    'application/vnd.ms-powerpoint': 'PPT',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'PPTX',
    'text/csv': 'CSV',
    'application/rtf': 'RTF',
    'application/vnd.oasis.opendocument.text': 'ODT',
    'application/vnd.oasis.opendocument.spreadsheet': 'ODS',
    'text/plain': 'TXT',
}
UA = ('Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 '
      '(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36')
HEADERS = {
    'User-Agent': UA,
    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8',
    'Accept-Language': 'en-US,en;q=0.9',
    'Accept-Encoding': 'gzip, deflate, br',
    'Connection': 'keep-alive',
    'Upgrade-Insecure-Requests': '1',
}
BLOCK_MARKERS = [
    'cf-challenge', 'just a moment', 'attention required', 'access denied',
    'enable javascript and cookies', 'verify you are human', 'checking your browser',
    'captcha', 'ddos protection', '403 forbidden',
]
CHARSETS = ['utf-8', 'cp1252', 'latin-1']


class Crawler:
    def __init__(self, args):
        self.args = args
        self.session = requests.Session()
        self.session.headers.update(HEADERS)
        self.base = urlparse(args.url)
        self.base_host = (self.base.hostname or '').lower()
        self.root_domain = self.base_host[4:] if self.base_host.startswith('www.') else self.base_host
        self.visited_pages = set()
        self.visited_docs = set()
        self.pages_seen = 0
        self.docs = []          # list of doc records
        self.page_errors = []
        self.blocked_hosts = set()
        self.robots_disallow = set()

    # ---------- helpers ----------
    def is_same_domain(self, url):
        host = (urlparse(url).hostname or '').lower()
        if not host:
            return False
        if self.args.include_other_domains:
            return True
        root = self.root_domain
        return host == root or host.endswith('.' + root)

    def normalize(self, url):
        p = urlparse(url)
        path = re.sub(r'/+', '/', p.path or '/')
        # drop common tracking/fragment junk
        q = re.sub(r'(&?(utm_[^=]+|gclid|fbclid|mc_cid|mc_eid)=[^&]*)', '', p.query).strip('&')
        return urlunparse((p.scheme.lower() or 'https', p.netloc.lower(), path, '', q, ''))

    def doc_type(self, url, content_type=''):
        path = unquote(urlparse(url).path).lower()
        for ext, kind in DOC_EXTENSIONS.items():
            if path.endswith(ext):
                return kind
        if content_type:
            ct = content_type.split(';')[0].strip().lower()
            return CONTENT_TYPE_MAP.get(ct)
        return None

    def polite(self):
        if self.args.delay > 0:
            time.sleep(self.args.delay)

    def looks_blocked(self, text, status):
        low = (text or '')[:4000].lower()
        if status in (403, 429) or status >= 500:
            return True
        return any(m in low for m in BLOCK_MARKERS)

    def fetch(self, url, stream=False, timeout=45):
        """GET with retries + backoff. Returns (requests.Response|None, error_str|None)."""
        self.polite()
        last_err = None
        for attempt in range(4):
            try:
                r = self.session.get(url, stream=stream, timeout=timeout,
                                     allow_redirects=True,
                                     verify=not self.args.insecure)
                body_preview = ''
                if not stream:
                    body_preview = r.text[:4000]
                if self.looks_blocked(body_preview, r.status_code):
                    # mark host; if a hard block (403/429), back off and retry
                    self.blocked_hosts.add(urlparse(url).hostname or '')
                    if r.status_code in (403, 429):
                        time.sleep(2 ** attempt * 2)
                        continue
                if r.status_code >= 400:
                    last_err = f'HTTP {r.status_code}'
                    time.sleep(2 ** attempt * 1.5)
                    continue
                return r, None
            except requests.exceptions.SSLError as e:
                last_err = f'SSL error: {e.__class__.__name__}'
                if not self.args.insecure:
                    break
                time.sleep(1.5 * (attempt + 1))
            except (requests.exceptions.ConnectionError, requests.exceptions.Timeout,
                    requests.exceptions.ChunkedEncodingError) as e:
                last_err = f'{e.__class__.__name__}'
                time.sleep(2 ** attempt * 1.5)
            except Exception as e:
                last_err = f'{e.__class__.__name__}: {e}'
                break
        return None, last_err or 'unreachable'

    # ---------- crawling ----------
    def load_robots(self):
        if not self.args.enforce_robots:
            return
        try:
            r, _ = self.fetch(urljoin(self.args.url, '/robots.txt'), timeout=15)
            if r:
                for line in r.text.splitlines():
                    m = re.match(r'^Disallow:\s*(\S+)', line.strip())
                    if m:
                        self.robots_disallow.add(m.group(1))
        except Exception:
            pass

    def robots_allowed(self, url):
        if not self.robots_disallow:
            return True
        path = urlparse(url).path
        return not any(path.startswith(d) for d in self.robots_disallow if d)

    def walk_json(self, text):
        """Recursively collect URL strings from a JSON API payload (SPA finders)."""
        try:
            data = json.loads(text)
        except Exception:
            return []
        out = []
        def rec(o):
            if isinstance(o, dict):
                for v in o.values():
                    rec(v)
            elif isinstance(o, list):
                for v in o:
                    rec(v)
            elif isinstance(o, str):
                s = o.strip()
                if s.startswith(('http://', 'https://')):
                    out.append(s)
                elif s.startswith('/') and len(s) > 2:
                    low = s.lower()
                    if low.endswith(tuple(DOC_EXTENSIONS)) or '/api/' in low:
                        out.append(s)
                    # content-page paths (no asset extension, >= 2 segments)
                    elif not low.endswith(('.css', '.js', '.png', '.jpg', '.jpeg', '.gif',
                                            '.svg', '.ico', '.woff', '.woff2', '.ttf', '.map',
                                            '.json', '.xml')) and s.count('/') >= 2:
                        out.append(s)
        rec(data)
        return out

    def crawl(self):
        self.load_robots()
        queue = [self.normalize(self.args.url)]
        queued = set(queue)
        for seed in self.args.seed_pages:
            u = self.normalize(urljoin(self.args.url, seed))
            if u not in queued:
                queue.append(u)
                queued.add(u)

        while queue and self.pages_seen < self.args.max_pages:
            url = queue.pop(0)
            if url in self.visited_pages or not self.is_same_domain(url):
                continue
            if len(queued) > 8000:  # safety cap against runaway link farms
                break
            if not self.robots_allowed(url):
                continue
            self.visited_pages.add(url)
            self.pages_seen += 1

            r, err = self.fetch(url)
            if err or r is None:
                self.page_errors.append({'url': url, 'error': err})
                continue

            ct = r.headers.get('Content-Type', '')
            # If the "page" is actually a document, grab it directly
            d = self.doc_type(url, ct)
            if d and 'html' not in ct:
                self.process_document(url, d, r.headers.get('Content-Length'))
                continue

            html = r.text
            # JSON API response (SPA finder apps): walk for document/page URLs
            if 'json' in ct or html.lstrip()[:1] in ('{', '['):
                for target in self.walk_json(html):
                    if not target.startswith('http'):
                        target = urljoin(url, target)
                    target = self.normalize(target)
                    if not target.startswith(('http://', 'https://')):
                        continue
                    dt = self.doc_type(target)
                    if dt and target not in self.visited_docs:
                        self.process_document(target, dt)
                    elif self.is_same_domain(target) and target not in queued:
                        queue.append(target)
                        queued.add(target)
                continue

            soup = BeautifulSoup(html, 'lxml')
            for a in soup.find_all('a', href=True):
                href = a['href'].strip()
                if href.startswith(('mailto:', 'tel:', 'javascript:', '#')):
                    continue
                target = self.normalize(urljoin(url, href))
                if not target.startswith(('http://', 'https://')):
                    continue
                dt = self.doc_type(target)
                if dt:
                    if target not in self.visited_docs:
                        self.process_document(target, dt)
                elif self.is_same_domain(target):
                    queue.append(target)

            # JS/SPA fallback: hunt for document URLs inside inline scripts & JSON
            for blob in re.findall(r'<script[^>]*>(.*?)</script>', html, re.S | re.I)[:8]:
                for m in re.finditer(r'https?://[^\s"\'\\]+\.(?:pdf|docx?|xlsx?|pptx?|csv|rtf|odt|ods)(?:[?#][^\s"\'\\]*)?', blob, re.I):
                    target = self.normalize(m.group(0))
                    if self.is_same_domain(target) and target not in self.visited_docs:
                        self.process_document(target, self.doc_type(target))

            # Queue JSON API endpoints hinted by the page (SPA finder apps)
            for api in re.findall(r'["\'](/api/[^"\']{2,120})["\']', html)[:20]:
                api_url = self.normalize(urljoin(url, api))
                if self.is_same_domain(api_url) and api_url not in queued:
                    queue.append(api_url)
                    queued.add(api_url)

        # write crawl stats
        stats = {
            'target_url': self.args.url,
            'pages_visited': self.pages_seen,
            'max_pages_cap': self.args.max_pages,
            'documents_found': len(self.docs),
            'blocked_hosts': sorted(self.blocked_hosts),
            'page_errors': self.page_errors[:50],
            'warnings': [],
        }
        if self.blocked_hosts:
            stats['warnings'].append(
                'Bot protection detected on: ' + ', '.join(sorted(self.blocked_hosts)) +
                '. If pages were blocked, retry with Firecrawl (JS rendering) as fallback.')
        if self.pages_seen >= self.args.max_pages:
            stats['warnings'].append('Hit max-pages cap; re-run with --max-pages higher for a deeper crawl.')
        with open(os.path.join(self.args.output_dir, 'crawl_stats.json'), 'w', encoding='utf-8') as f:
            json.dump(stats, f, ensure_ascii=False, indent=2)
        return stats

    # ---------- documents ----------
    def process_document(self, url, doc_type, known_size=None):
        self.visited_docs.add(url)
        fname = self.safe_filename(url, doc_type)
        doc = {
            'url': url, 'type': doc_type, 'filename': fname,
            'size_bytes': None, 'http_status': None,
            'download_ok': False, 'download_error': None,
            'text_path': None, 'language': None, 'lang_conf': None,
            'title_hint': self.title_hint(url), 'char_count': 0,
            'scanned_pdf': False, 'translated': False,
        }
        self.docs.append(doc)
        try:
            r, err = self.fetch(url, stream=True, timeout=90)
            if err or r is None:
                doc['download_error'] = err or 'unreachable'
                return
            doc['http_status'] = r.status_code
            ctype = r.headers.get('Content-Type', '')
            detected = self.doc_type(url, ctype)
            if detected and detected != doc_type:
                doc['type'] = detected
                doc_type = detected
            length = int(r.headers.get('Content-Length') or 0)
            if length > self.args.max_file_mb * 1024 * 1024:
                doc['download_error'] = f'skipped: > {self.args.max_file_mb} MB'
                r.close()
                return
            doc['size_bytes'] = length or None
            path = os.path.join(self.args.output_dir, 'docs', fname)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            wrote = 0
            with open(path, 'wb') as f:
                for chunk in r.iter_content(chunk_size=65536):
                    f.write(chunk)
                    wrote += len(chunk)
                    if wrote > self.args.max_file_mb * 1024 * 1024:
                        doc['download_error'] = f'stopped at {self.args.max_file_mb} MB cap'
                        break
            r.close()
            doc['size_bytes'] = wrote
            doc['download_ok'] = True
            text = self.extract_text(path, doc_type)
            if text is None:
                doc['download_error'] = 'text extraction unsupported for this format'
            elif text.strip():
                doc['char_count'] = len(text)
                if doc_type == 'PDF' and len(text.strip()) < 40:
                    doc['scanned_pdf'] = True
                lang, conf = self.detect_language(text)
                doc['language'] = lang
                doc['lang_conf'] = round(float(conf), 3) if conf is not None else None
                text = text[:self.args.max_text_chars]
                tpath = os.path.join(self.args.output_dir, 'text', fname + '.txt')
                os.makedirs(os.path.dirname(tpath), exist_ok=True)
                with open(tpath, 'w', encoding='utf-8') as f:
                    f.write(text)
                doc['text_path'] = os.path.relpath(tpath, self.args.output_dir)
                if self.args.translate and lang and lang != 'en':
                    doc['translated'] = self.try_translate(text, lang)
            elif doc_type == 'PDF' and self.args.enable_ocr:
                doc['scanned_pdf'] = True
                text = self.ocr_pdf(path)
                if text:
                    doc['char_count'] = len(text)
                    doc['ocr_applied'] = True
                    lang, conf = self.detect_language(text)
                    doc['language'] = lang
                    doc['lang_conf'] = round(float(conf), 3) if conf is not None else None
                    text = text[:self.args.max_text_chars]
                    tpath = os.path.join(self.args.output_dir, 'text', fname + '.txt')
                    os.makedirs(os.path.dirname(tpath), exist_ok=True)
                    with open(tpath, 'w', encoding='utf-8') as f:
                        f.write(text)
                    doc['text_path'] = os.path.relpath(tpath, self.args.output_dir)
                    if self.args.translate and lang and lang != 'en':
                        doc['translated'] = self.try_translate(text, lang)
                else:
                    doc['download_error'] = 'scanned PDF — OCR failed or unavailable'
            elif doc_type == 'PDF':
                doc['scanned_pdf'] = True
                doc['download_error'] = 'scanned PDF — enable OCR to extract text'
        except Exception as e:
            doc['download_error'] = f'{e.__class__.__name__}: {e}'

    def extract_text(self, path, doc_type):
        try:
            if doc_type == 'PDF' and HAS_FITZ:
                with fitz.open(path) as doc:
                    return '\n'.join(p.get_text() for p in doc)
            if doc_type == 'DOCX' and HAS_DOCX:
                d = DocxDocument(path)
                return '\n'.join(p.text for p in d.paragraphs) + '\n' + \
                       '\n'.join(c.text for t in d.tables for row in t.rows for c in row.cells)
            if doc_type == 'XLSX' and HAS_XLSX:
                wb = load_workbook(path, read_only=True, data_only=True)
                out = []
                for ws in wb.worksheets[:3]:
                    for row in ws.iter_rows(values_only=True):
                        if any(v is not None for v in row):
                            out.append(' | '.join(str(v) for v in row if v is not None))
                wb.close()
                return '\n'.join(out)
            if doc_type == 'XLS' and HAS_XLS:
                wb = xlrd.open_workbook(path)
                out = []
                for ws in wb.sheets()[:3]:
                    for r in range(min(ws.nrows, 200)):
                        out.append(' | '.join(str(ws.cell_value(r, c)) for c in range(ws.ncols)))
                return '\n'.join(out)
            if doc_type == 'PPTX' and HAS_PPTX:
                prs = Presentation(path)
                out = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            out.append(shape.text)
                return '\n'.join(out)
            if doc_type in ('CSV', 'TXT', 'RTF', 'ODT', 'ODS', 'DOC', 'PPT'):
                raw = open(path, 'rb').read()
                if doc_type in ('ODT', 'ODS'):
                    with zipfile.ZipFile(path) as z:
                        xml = z.read('content.xml').decode('utf-8', 'ignore')
                    return re.sub(r'<[^>]+>', ' ', xml)
                if doc_type == 'RTF':
                    return re.sub(r'\\[a-z]+-?\d* ?', ' ', raw.decode('cp1252', 'ignore'))
                text = None
                for enc in CHARSETS:
                    try:
                        text = raw.decode(enc)
                        break
                    except (UnicodeDecodeError, ValueError):
                        continue
                return text or ''
            return None
        except Exception:
            return None

    def detect_language(self, text):
        if not HAS_LANGID:
            try:
                from langdetect import detect, detect_langs
                ls = detect_langs(text[:2000])
                return str(ls[0].lang), ls[0].prob
            except Exception:
                return None, 0.0
        try:
            lang, prob = _lid.classify(text[:4000])
            return lang, prob
        except Exception:
            return None, 0.0

    def try_translate(self, text, lang):
        if not HAS_TRANSLATOR:
            return False
        try:
            src = lang if lang != 'zh' else 'zh-CN'
            chunk = text[:4000]
            out = GoogleTranslator(source='auto', target='en').translate(chunk)
            if out and out.strip():
                tpath = os.path.join(self.args.output_dir, 'text', 'translated_' + os.path.basename(self.docs[-1]['filename']) + '.txt')
                with open(tpath, 'w', encoding='utf-8') as f:
                    f.write(out)
                return True
        except Exception:
            return False
        return False

    def ocr_pdf(self, path):
        """OCR a scanned/image-based PDF using tesserocr + pdf2image."""
        if not HAS_PDF2IMAGE or not HAS_TESSEROCR:
            return None
        try:
            images = convert_from_path(path, dpi=300, fmt='png', grayscale=True)
            if not images:
                return None
            texts = []
            api = PyTessBaseAPI(psm=PSM.AUTO, path=_TESSDATA)
            for img in images:
                api.SetImage(img)
                chunk = api.GetUTF8Text().strip()
                if chunk:
                    texts.append(chunk)
            api.End()
            return '\n'.join(texts)
        except Exception as e:
            return None

    # ---------- misc ----------
    @staticmethod
    def title_hint(url):
        path = unquote(urlparse(url).path)
        name = os.path.basename(path).rsplit('.', 1)[0]
        name = re.sub(r'[-_]+', ' ', name)
        return ' '.join(w.capitalize() for w in name.split()) if name else url

    @staticmethod
    def safe_filename(url, doc_type):
        path = unquote(urlparse(url).path)
        name = os.path.basename(path)
        if not name or '.' not in name:
            name = 'doc_' + str(abs(hash(url)) % 10 ** 8)
        name = re.sub(r'[^\w.\-]', '_', name)
        if not name.lower().endswith(tuple(DOC_EXTENSIONS)):
            name += DOC_EXTENSIONS.get('.' + name.split('.')[-1].lower(), '.' + doc_type.lower())
        return name[:180]


def main():
    ap = argparse.ArgumentParser(description='Ministry document crawler')
    ap.add_argument('--url', required=True)
    ap.add_argument('--output-dir', default='./crawl_out')
    ap.add_argument('--max-pages', type=int, default=300)
    ap.add_argument('--delay', type=float, default=0.8)
    ap.add_argument('--max-file-mb', type=int, default=150)
    ap.add_argument('--max-text-chars', type=int, default=20000)
    ap.add_argument('--include-other-domains', action='store_true')
    ap.add_argument('--enforce-robots', action='store_true')
    ap.add_argument('--translate', action='store_true')
    ap.add_argument('--enable-ocr', action='store_true')
    ap.add_argument('--insecure', action='store_true')
    ap.add_argument('--seed-pages', nargs='*', default=[])
    args = ap.parse_args()

    os.makedirs(args.output_dir, exist_ok=True)
    os.makedirs(os.path.join(args.output_dir, 'docs'), exist_ok=True)
    os.makedirs(os.path.join(args.output_dir, 'text'), exist_ok=True)

    crawler = Crawler(args)
    print(f'[crawler] target={args.url} max_pages={args.max_pages} delay={args.delay}s', flush=True)
    t0 = time.time()
    stats = crawler.crawl()
    elapsed = round(time.time() - t0, 1)

    # index.json
    with open(os.path.join(args.output_dir, 'index.json'), 'w', encoding='utf-8') as f:
        json.dump(crawler.docs, f, ensure_ascii=False, indent=2)

    # index.csv (quick view)
    with open(os.path.join(args.output_dir, 'index.csv'), 'w', newline='', encoding='utf-8') as f:
        w = csv.DictWriter(f, fieldnames=['filename', 'type', 'url', 'size_bytes', 'language',
                                          'download_ok', 'download_error', 'text_path',
                                          'title_hint', 'char_count', 'scanned_pdf'])
        w.writeheader()
        for d in crawler.docs:
            w.writerow({k: d.get(k) for k in w.fieldnames})

    ok = sum(1 for d in crawler.docs if d['download_ok'])
    txt = sum(1 for d in crawler.docs if d['text_path'])
    langs = {}
    for d in crawler.docs:
        if d['language']:
            langs[d['language']] = langs.get(d['language'], 0) + 1

    print(f'[crawler] DONE in {elapsed}s | pages={stats["pages_visited"]} '
          f'docs_found={len(crawler.docs)} downloaded_ok={ok} with_text={txt} '
          f'languages={json.dumps(langs, ensure_ascii=False)}')
    print(f'[crawler] output_dir={os.path.abspath(args.output_dir)}')
    if crawler.blocked_hosts:
        print(f'[crawler] WARNING blocked by bot-protection on: {", ".join(sorted(crawler.blocked_hosts))} '
              '- consider Firecrawl fallback for JS/blocked pages')
    if crawler.page_errors:
        print(f'[crawler] WARNING {len(crawler.page_errors)} page errors (see crawl_stats.json)')


if __name__ == '__main__':
    main()
